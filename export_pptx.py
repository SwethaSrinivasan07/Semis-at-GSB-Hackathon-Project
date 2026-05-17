"""Screenshot each slide at 1280x720 and pack into a .pptx file."""
import asyncio, os, pathlib
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

SLIDES_URL = "http://localhost:8888/slides.html"
OUT_DIR = pathlib.Path(__file__).parent / "slide_images"
OUT_PPTX = pathlib.Path(__file__).parent / "Catena.pptx"
W, H = 1280, 720

async def screenshot_slides():
    OUT_DIR.mkdir(exist_ok=True)
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": W, "height": H})
        await page.goto(SLIDES_URL, wait_until="networkidle")
        await page.wait_for_timeout(1500)  # let fonts settle

        slides = await page.query_selector_all(".slide")
        print(f"Found {len(slides)} slides")
        paths = []
        for i, slide in enumerate(slides):
            path = OUT_DIR / f"slide_{i+1:02d}.png"
            await slide.screenshot(path=str(path))
            print(f"  Saved {path.name}")
            paths.append(path)
        await browser.close()
        return paths

def build_pptx(image_paths):
    prs = Presentation()
    prs.slide_width  = Emu(W * 914400 // 96)   # 1280px @ 96dpi → EMU
    prs.slide_height = Emu(H * 914400 // 96)    # 720px  @ 96dpi → EMU

    blank_layout = prs.slide_layouts[6]  # blank
    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path), 0, 0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    prs.save(str(OUT_PPTX))
    print(f"\nSaved → {OUT_PPTX}")

async def main():
    paths = await screenshot_slides()
    build_pptx(paths)

asyncio.run(main())
