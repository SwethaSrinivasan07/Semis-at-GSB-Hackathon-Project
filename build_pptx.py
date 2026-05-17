"""
Build Catena_editable.pptx — fully editable, native python-pptx shapes and text.
Run: /opt/homebrew/bin/python3 build_pptx.py
"""
import pathlib
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Slide dimensions: 1280×720 px @ 96 dpi ─────────────────────
PX = 9525          # 1 px → EMU
SW = 1280 * PX
SH =  720 * PX

# ── Design tokens ───────────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x3A, 0x5F)
BLUE    = RGBColor(0x25, 0x63, 0xEB)
BLUE_M  = RGBColor(0x60, 0xA5, 0xFA)
BLUE_L  = RGBColor(0xBF, 0xDB, 0xFE)
BLUE_BG = RGBColor(0xF0, 0xF7, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
RED     = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
GREEN   = RGBColor(0x05, 0x96, 0x69)
DARK_BG = RGBColor(0x1E, 0x2D, 0x4A)
BLUE_93 = RGBColor(0x93, 0xC5, 0xFD)

SERIF = "Playfair Display"
SANS  = "Plus Jakarta Sans"

OUT = pathlib.Path(__file__).parent / "Catena_editable.pptx"

# ── Primitive helpers ────────────────────────────────────────────
def px(n):
    return int(n * PX)

def box(slide, x, y, w, h, fill, border=None, border_pt=0.75):
    """Plain filled rectangle. No XML tricks."""
    s = slide.shapes.add_shape(1, px(x), px(y), px(w), px(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s

def circ(slide, x, y, w, h, fill):
    """Oval / circle shape."""
    s = slide.shapes.add_shape(9, px(x), px(y), px(w), px(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, x, y, w, h, text, font, size, color,
        bold=False, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Single-paragraph textbox."""
    shape = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return shape

def txt2(slide, x, y, w, h, t1, c1, t2, c2, font, size, bold=True):
    """Two-color text in one paragraph."""
    shape = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    for text, color in [(t1, c1), (t2, c2)]:
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return shape

def send_back(slide, shape):
    t = slide.shapes._spTree
    t.remove(shape._element)
    t.insert(2, shape._element)

# ── Shared chrome ────────────────────────────────────────────────
def bg(slide, color=BLUE_BG):
    s = box(slide, 0, 0, 1280, 720, color)
    send_back(slide, s)

def blobs(slide, light=BLUE_L, mid=BLUE_M):
    circ(slide,    0,   0, 280, 280, light)   # TL (clipped)
    circ(slide, 1090, -30, 210, 210, mid)     # TR
    circ(slide, 1090, 510, 250, 250, light)   # BR
    circ(slide,    0, 570, 170, 170, mid)     # BL

def rule(slide, color=BLUE):
    box(slide, 0, 714, 1280, 6, color)

def tag(slide, label):
    txt(slide, 980, 20, 280, 22, label, SANS, 8, BLUE,
        bold=True, align=PP_ALIGN.RIGHT)

def eyebrow(slide, x, y, label, color=BLUE):
    txt(slide, x, y, 700, 20, label.upper(), SANS, 8, color, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════
def s1(slide):
    bg(slide)
    blobs(slide)
    rule(slide)
    tag(slide, "SEMIS @ GSB · MAY 2026")

    # Badge
    box(slide, 100, 168, 152, 26, BLUE)
    txt(slide, 106, 172, 148, 20, "Hackathon Project",
        SANS, 9, WHITE, bold=True)

    # Title: "Cat" (navy) + "ena" (blue)
    txt2(slide, 96, 200, 680, 105, "Cat", NAVY, "ena", BLUE,
         SERIF, 80, bold=True)

    # Subtitle
    txt(slide, 100, 310, 500, 58,
        "Design-time supply chain intelligence for photonics OEMs.",
        SANS, 17, MUTED, wrap=True)

    # Team block
    txt(slide, 100, 406, 240, 20, "Swetha, Jed, Parul, Josh",
        SANS, 12, NAVY, bold=True)
    txt(slide, 100, 428, 240, 16, "Team Catena", SANS, 10, MUTED)

    # Thin vertical rule between team / event
    box(slide, 360, 404, 1, 38, BLUE_L)

    # Event
    txt(slide, 372, 406, 260, 20, "Semis @ GSB Hackathon",
        SANS, 12, BLUE, bold=True)
    txt(slide, 372, 428, 260, 16, "Stanford GSB · May 2026",
        SANS, 10, MUTED)

    # ── Stats card (right side) ─────────────────────────────────
    box(slide, 830, 148, 370, 326, WHITE, border=BLUE_L)

    txt(slide, 856, 168, 318, 16,
        "THE NUMBERS THAT MATTER", SANS, 8, MUTED, bold=True)

    # Stat 1 — 52 weeks
    txt(slide, 856, 192, 318, 46, "52 weeks",
        SERIF, 32, RED, bold=True)
    txt(slide, 856, 238, 318, 32,
        "InP photonic chip lead time — discovered after tape-out",
        SANS, 9, MUTED, wrap=True)
    box(slide, 856, 276, 318, 1, BLUE_L)

    # Stat 2 — 11–13%
    txt(slide, 856, 284, 318, 46, "11–13%",
        SERIF, 32, BLUE, bold=True)
    txt(slide, 856, 330, 318, 24,
        "margin Avnet and Arrow take on a $100B market",
        SANS, 9, MUTED, wrap=True)
    box(slide, 856, 360, 318, 1, BLUE_L)

    # Stat 3 — 20%
    txt(slide, 856, 368, 318, 46, "20%",
        SERIF, 32, AMBER, bold=True)
    txt(slide, 856, 414, 318, 32,
        "Arrow revenue decline in 2024 — disruption already started",
        SANS, 9, MUTED, wrap=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ═══════════════════════════════════════════════════════════════
def s2(slide):
    bg(slide)
    blobs(slide)
    rule(slide)
    tag(slide, "THE PROBLEM")
    eyebrow(slide, 80, 46, "A design-time information problem")

    txt(slide, 80, 66, 900, 138,
        "Engineers lock in supply chains 12 months\nbefore procurement ever looks.",
        SERIF, 30, NAVY, bold=True, wrap=True)

    # Row 1 cards  (y=220)
    box(slide, 80,  220, 556, 128, WHITE, border=BLUE_L)
    txt(slide, 102, 232, 160,  52, "3",    SERIF, 40, RED,   bold=True)
    txt(slide, 102, 280, 512,  54,
        "InP wafer fabs globally. Two in geopolitically sensitive regions. Zero alternatives for coherent photonic chips.",
        SANS, 9, MUTED, wrap=True)

    box(slide, 650, 220, 550, 128, WHITE, border=BLUE_L)
    txt(slide, 672, 232, 310,  52, "5–7 yrs", SERIF, 40, BLUE, bold=True)
    txt(slide, 672, 280, 506,  54,
        "One component choice at design time locks in a supply chain for the life of the program.",
        SANS, 9, MUTED, wrap=True)

    # Row 2 cards  (y=362)
    box(slide, 80,  362, 556, 128, WHITE, border=BLUE_L)
    txt(slide, 102, 374, 220,  52, "60%",  SERIF, 40, AMBER, bold=True)
    txt(slide, 102, 422, 512,  54,
        "of Coherent's transceiver capacity allocated to hyperscalers. OEMs compete for the rest.",
        SANS, 9, MUTED, wrap=True)

    box(slide, 650, 362, 550, 128, NAVY)
    txt(slide, 672, 374, 210,  52, "$0",   SERIF, 40, WHITE, bold=True)
    txt(slide, 672, 422, 506,  54,
        "spent on supply chain tooling at design time. SiliconExpert is owned by Arrow — the distributor with a conflict of interest.",
        SANS, 9, BLUE_93, wrap=True)

    # Timeline bar  (y=506)
    box(slide, 80, 506, 1120, 90, WHITE, border=BLUE_L)
    txt(slide, 102, 518, 700, 14,
        "WHERE SUPPLY CHAIN RISK IS CREATED VS. DISCOVERED",
        SANS, 8, MUTED, bold=True)
    box(slide, 102, 538, 530, 30, BLUE)
    txt(slide, 110, 544, 520, 18,
        "Design time — risk created here", SANS, 9, WHITE, bold=True)
    box(slide, 635, 538, 266, 30, BLUE_M)
    txt(slide, 643, 544, 250, 18, "Prototyping", SANS, 9, NAVY, bold=True)
    box(slide, 904, 538, 274, 30, RED)
    txt(slide, 912, 544, 258, 18, "Procurement scramble", SANS, 9, WHITE, bold=True)
    txt(slide, 102, 574, 1076, 14,
        "By the time procurement discovers a 52-week lead time, the design is already locked. Nobody has a tool for the left side of this bar.",
        SANS, 9, MUTED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Product
# ═══════════════════════════════════════════════════════════════
def s3(slide):
    bg(slide)
    blobs(slide)
    rule(slide)
    tag(slide, "THE PRODUCT")

    txt(slide, 80, 38, 160, 18, "Catena", SANS, 10, BLUE, bold=True)
    txt(slide, 80, 58, 560, 130,
        "Upload a BOM.\nKnow every risk\nbefore tape-out.",
        SERIF, 30, NAVY, bold=True, wrap=True)

    # 4-step list  (starts at y=198)
    steps = [
        ("1", "Upload your BOM",
         "CSV or Excel. Claude parses it instantly — no templates required."),
        ("2", "Risk scored across 6 dimensions",
         "Availability, lead time, cost, lifecycle, geopolitical exposure, vendor concentration. Every part: Red / Yellow / Green."),
        ("3", "Ranked substitutes, ready to act on",
         "Drop-in, minor rework, or redesign — with compatibility grade, qualification timeline, and why each one works."),
        ("4", "Annotated BOM exported",
         "Color-coded Excel. Engineer sends to procurement. Procurement sends to leadership. Same picture in minutes."),
    ]
    sy = 196
    for num, title, desc in steps:
        circ(slide, 80, sy, 26, 26, BLUE)
        txt(slide, 80, sy + 4, 26, 18, num, SANS, 9, WHITE,
            bold=True, align=PP_ALIGN.CENTER)
        txt(slide, 116, sy,      450, 20, title, SANS, 11, NAVY, bold=True)
        txt(slide, 116, sy + 20, 450, 36, desc,  SANS,  9, MUTED, wrap=True)
        sy += 80

    # ── Right panel ─────────────────────────────────────────────
    # Market card
    box(slide, 636, 148, 564, 96, NAVY)
    txt(slide, 658, 160, 300, 14, "Target market", SANS, 8, BLUE_93, bold=True)
    txt(slide, 658, 174, 500, 50, "$10–15B", SERIF, 34, WHITE, bold=True)
    txt(slide, 658, 222, 520, 18,
        "photonics transceiver market/yr. 2% capture = $200–300M. Pure orchestration.",
        SANS, 9, BLUE_93, wrap=True)

    # Competitor table
    box(slide, 636, 258, 564, 174, WHITE, border=BLUE_L)
    txt(slide, 658, 272, 300, 14, "vs. existing tools", SANS, 8, MUTED, bold=True)

    rows = [
        ("SiliconExpert",    "owned by Arrow",     RED),
        ("Octopart",         "no AVL, no AI",       MUTED),
        ("Avnet / Arrow reps","2-week turnaround",   MUTED),
        ("Catena",           "design-time, AI-native", GREEN),
    ]
    ry = 292
    for name, label, color in rows:
        bold_row = name == "Catena"
        txt(slide, 658, ry, 270, 18, name,  SANS, 10, NAVY,  bold=bold_row)
        txt(slide, 852, ry, 326, 18, label, SANS, 10, color, bold=bold_row,
            align=PP_ALIGN.RIGHT)
        ry += 34

    # Pricing card
    box(slide, 636, 446, 564, 104, BLUE_BG, border=BLUE_L)
    txt(slide, 658, 460, 400, 14, "Phase 1 revenue model", SANS, 8, MUTED, bold=True)
    txt(slide, 658, 476, 500, 44, "$50–200K / year", SERIF, 28, NAVY, bold=True)
    txt(slide, 658, 518, 520, 26,
        "SaaS per OEM. Phase 2: 1.5–2.5% on direct OEM-to-manufacturer transactions.",
        SANS, 9, MUTED, wrap=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Demo
# ═══════════════════════════════════════════════════════════════
def s4(slide):
    bg(slide, DARK_BG)
    circ(slide,    0,   0, 220, 220, RGBColor(0x1D, 0x4E, 0xD8))
    circ(slide, 1100, 520, 220, 220, RGBColor(0x3B, 0x82, 0xF6))
    rule(slide, BLUE_M)

    txt(slide, 540, 268, 200, 72, "▶", SANS, 44, BLUE_M,
        align=PP_ALIGN.CENTER)
    txt(slide, 280, 344, 720, 76, "Live demo", SERIF, 52, WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    txt(slide, 280, 428, 720, 30, "Video embed goes here", SANS, 16, BLUE_M,
        align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Build + Learnings
# ═══════════════════════════════════════════════════════════════
def s5(slide):
    bg(slide)
    blobs(slide)
    rule(slide)
    tag(slide, "HOW WE BUILT IT")
    eyebrow(slide, 80, 40, "The build")

    txt(slide, 80, 58, 800, 96,
        "Four workstreams, one weekend.",
        SERIF, 30, NAVY, bold=True, wrap=True)

    # ── Left: What we built ─────────────────────────────────────
    txt(slide, 80, 168, 400, 18, "What we built",
        SANS, 9, MUTED, bold=True)

    # Streamlit card
    box(slide, 80, 188, 496, 108, WHITE, border=BLUE_L)
    txt(slide, 100, 200, 456, 20,
        "Streamlit dashboard + FastAPI backend",
        SANS, 11, NAVY, bold=True)
    txt(slide, 100, 222, 456, 44,
        "Full-stack app: BOM upload, risk dashboard, substitute cards, annotated Excel export. Deployed from Git with four parallel feature branches.",
        SANS, 9, MUTED, wrap=True)

    # Tech tags
    tags = [("Python", BLUE), ("Streamlit", BLUE_M),
            ("FastAPI", NAVY), ("Claude API", BLUE)]
    tx = 100
    ty = 270
    for label, col in tags:
        tw = len(label) * 7 + 18
        box(slide, tx, ty, tw, 20, col)
        txt(slide, tx + 4, ty + 3, tw - 6, 14, label,
            SANS, 8, WHITE, bold=True)
        tx += tw + 8

    # Claude card
    box(slide, 80, 308, 496, 86, BLUE)
    txt(slide, 100, 320, 456, 20,
        "AI layer: Claude as the intelligence core",
        SANS, 11, WHITE, bold=True)
    txt(slide, 100, 342, 456, 46,
        "BOM parsing, 6-dimension risk scoring, photonics-specific substitution reasoning, and plain-English VP-level risk briefings. Powered by claude-sonnet-4-6.",
        SANS, 9, BLUE_L, wrap=True)

    # ── Right: What we learned ──────────────────────────────────
    txt(slide, 626, 168, 600, 18, "What we learned",
        SANS, 9, MUTED, bold=True)

    learnings = [
        ("The problem is deeper than procurement",
         "Every engineer we talked to confirmed: the issue is at design time, not sourcing time. No one has built a tool that lives there."),
        ("AVL awareness is the real moat",
         "Any tool can flag a long lead time. Knowing which alternates are on a customer's Approved Vendor List — and which need 6–18 months of qualification — is where the value actually lives."),
        ("Nokia / Infinera is the wedge",
         "$6.65B merger, two BOMs merging, two AVLs conflicting. Maximum supply chain chaos. Maximum willingness to pay."),
    ]
    ly = 188
    for title, desc in learnings:
        box(slide, 626, ly, 574, 96, WHITE, border=BLUE_L)
        txt(slide, 646, ly + 12, 534, 20, title, SANS, 11, NAVY, bold=True)
        txt(slide, 646, ly + 34, 534, 52, desc,  SANS,  9, MUTED, wrap=True)
        ly += 106


# ═══════════════════════════════════════════════════════════════
# Assemble
# ═══════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)
    blank = prs.slide_layouts[6]

    for fn, label in [(s1, "Cover"), (s2, "Problem"), (s3, "Product"),
                      (s4, "Demo"),  (s5, "Build + Learnings")]:
        slide = prs.slides.add_slide(blank)
        fn(slide)
        print(f"  Built: {label}")

    prs.save(str(OUT))
    print(f"\nSaved → {OUT}")

build()
